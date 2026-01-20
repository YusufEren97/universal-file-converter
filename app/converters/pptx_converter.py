"""
PPTX Converter
Converts PowerPoint files to PDF and images.
Uses PyMuPDF for PDF to image conversion (no Poppler needed).
"""
import os
import asyncio
import subprocess
import shutil


async def convert_pptx(input_path: str, output_dir: str, target_format: str) -> dict:
    """PPTX converter supporting PDF and image outputs."""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_pptx, input_path, output_dir, target_format)


def _process_pptx(input_path: str, output_dir: str, target_format: str) -> dict:
    try:
        filename = os.path.basename(input_path)
        name, ext = os.path.splitext(filename)
        output_filename = f"{name}.{target_format}"
        output_path = os.path.join(output_dir, output_filename)

        if target_format == 'pdf':
            return _pptx_to_pdf(input_path, output_path, output_filename)
        elif target_format in ['png', 'jpg', 'jpeg']:
            return _pptx_to_images(input_path, output_dir, target_format, name)
        elif target_format == 'txt':
            return _pptx_to_txt(input_path, output_path, output_filename)
        else:
            return {"success": False, "error": f"Unsupported target format for PPTX: {target_format}"}

    except Exception as e:
        return {"success": False, "error": f"PPTX conversion error: {str(e)}"}


def _pptx_to_pdf(input_path: str, output_path: str, output_filename: str) -> dict:
    """Convert PPTX to PDF using LibreOffice or PowerPoint COM."""
    
    # Method 1: Try LibreOffice (cross-platform)
    try:
        libreoffice_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            "soffice",
            "/usr/bin/soffice",
            "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        ]
        
        soffice_path = None
        for path in libreoffice_paths:
            if os.path.exists(path) or shutil.which(path):
                soffice_path = path
                break
        
        if soffice_path:
            out_dir = os.path.dirname(output_path)
            result = subprocess.run([
                soffice_path,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", out_dir,
                input_path
            ], capture_output=True, text=True, timeout=180)
            
            expected_output = os.path.join(out_dir, os.path.splitext(os.path.basename(input_path))[0] + ".pdf")
            if os.path.exists(expected_output):
                if expected_output != output_path:
                    shutil.move(expected_output, output_path)
                return {"success": True, "output_path": output_path, "filename": output_filename}
    except Exception as e:
        print(f"[PPTX→PDF] LibreOffice failed: {e}")
    
    # Method 2: Try PowerPoint COM (Windows only)
    try:
        import platform
        if platform.system() == 'Windows':
            import comtypes.client
            
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = 1
            
            presentation = powerpoint.Presentations.Open(os.path.abspath(input_path))
            presentation.SaveAs(os.path.abspath(output_path), 32)  # 32 = ppSaveAsPDF
            presentation.Close()
            powerpoint.Quit()
            
            if os.path.exists(output_path):
                return {"success": True, "output_path": output_path, "filename": output_filename}
    except Exception as e:
        print(f"[PPTX→PDF] PowerPoint COM failed: {e}")
    
    return {
        "success": False, 
        "error": "PDF dönüşümü için LibreOffice veya PowerPoint gerekli. Lütfen birini yükleyin."
    }


def _pptx_to_images(input_path: str, output_dir: str, target_format: str, name: str) -> dict:
    """Convert PPTX to images using PyMuPDF (no Poppler needed)."""
    import zipfile
    
    # First convert to PDF
    pdf_path = os.path.join(output_dir, f"{name}_temp.pdf")
    pdf_result = _pptx_to_pdf(input_path, pdf_path, f"{name}.pdf")
    
    if not pdf_result["success"]:
        # Fallback: Extract embedded images from PPTX directly
        try:
            images_extracted = []
            with zipfile.ZipFile(input_path, 'r') as pptx:
                for i, item in enumerate(pptx.namelist()):
                    if item.startswith('ppt/media/') and any(item.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif']):
                        img_data = pptx.read(item)
                        original_ext = os.path.splitext(item)[1]
                        img_filename = f"{name}_image_{i+1}{original_ext}"
                        img_path = os.path.join(output_dir, img_filename)
                        with open(img_path, 'wb') as f:
                            f.write(img_data)
                        images_extracted.append(img_filename)
            
            if images_extracted:
                return {
                    "success": True, 
                    "output_path": os.path.join(output_dir, images_extracted[0]), 
                    "filename": images_extracted[0],
                    "note": f"Sunumdan {len(images_extracted)} resim çıkarıldı"
                }
            return {"success": False, "error": "Sunumda resim bulunamadı ve PDF dönüşümü başarısız"}
        except Exception as e:
            return {"success": False, "error": f"Resim çıkarma hatası: {str(e)}"}
    
    # Convert PDF to images using PyMuPDF (fitz) - NO POPPLER NEEDED!
    try:
        import fitz  # PyMuPDF
        
        pdf_doc = fitz.open(pdf_path)
        output_files = []
        
        for page_num in range(len(pdf_doc)):
            page = pdf_doc[page_num]
            # Higher resolution for better quality
            mat = fitz.Matrix(2.0, 2.0)  # 2x zoom for 144 DPI
            pix = page.get_pixmap(matrix=mat)
            
            img_filename = f"{name}_slide_{page_num + 1}.{target_format}"
            img_path = os.path.join(output_dir, img_filename)
            
            if target_format.lower() in ['jpg', 'jpeg']:
                pix.save(img_path, output="jpeg", jpg_quality=95)
            else:
                pix.save(img_path)
            
            output_files.append(img_filename)
        
        pdf_doc.close()
        
        # Clean up temp PDF
        if os.path.exists(pdf_path):
            os.remove(pdf_path)
        
        if output_files:
            return {
                "success": True, 
                "output_path": os.path.join(output_dir, output_files[0]), 
                "filename": output_files[0],
                "all_files": output_files,
                "note": f"{len(output_files)} slayt resmi oluşturuldu"
            }
        
        return {"success": False, "error": "Sayfa dönüştürülemedi"}
        
    except ImportError:
        return {"success": False, "error": "PyMuPDF yüklü değil. 'pip install PyMuPDF' çalıştırın."}
    except Exception as e:
        # Clean up on error
        if os.path.exists(pdf_path):
            os.remove(pdf_path)
        return {"success": False, "error": f"PDF→resim dönüşüm hatası: {str(e)}"}


def _pptx_to_txt(input_path: str, output_path: str, output_filename: str) -> dict:
    """Extract text from PPTX."""
    try:
        from pptx import Presentation
    except ImportError:
        return {"success": False, "error": "python-pptx yüklü değil. 'pip install python-pptx' çalıştırın."}
    
    prs = Presentation(input_path)
    text_content = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        text_content.append(f"=== Slayt {slide_num} ===\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_content.append(shape.text)
        text_content.append("\n")
    
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write("\n".join(text_content))
    
    return {"success": True, "output_path": output_path, "filename": output_filename}
